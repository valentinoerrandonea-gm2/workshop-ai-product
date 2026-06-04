#!/usr/bin/env python3
"""Genera vision-makerlab-academy.pptx desde el guion de slides en textos/."""
import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

BASE = Path(__file__).resolve().parents[1]
SRC = BASE / "textos" / "vision-makerlab-academy.txt"
OUT = BASE.parent / "caso-makerlab" / "material-cliente" / "vision-makerlab-academy.pptx"


def parse_slides(text: str) -> list[tuple[str, list[str]]]:
    slides = []
    for bloque in re.split(r"^## SLIDE \d+: ", text, flags=re.M)[1:]:
        lineas = [l.strip() for l in bloque.strip().splitlines() if l.strip()]
        titulo = lineas[0]
        bullets = [l[2:] for l in lineas[1:] if l.startswith("- ")]
        slides.append((titulo, bullets))
    return slides


def main() -> None:
    prs = Presentation()
    layout = prs.slide_layouts[1]  # título + contenido

    for titulo, bullets in parse_slides(SRC.read_text()):
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = titulo
        body = slide.placeholders[1].text_frame
        body.clear()
        for i, b in enumerate(bullets):
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            p.text = b
            p.font.size = Pt(20)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"OK → {OUT}")


if __name__ == "__main__":
    main()
